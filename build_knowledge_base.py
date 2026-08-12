"""
BUILD KNOWLEDGE BASE — run this FIRST, only once (or whenever you add new files)
---------------------------------------------------------------------------------
This script:
1. Reads every supported file inside the "Stats notes" folder:
     - PDF (.pdf)
     - PowerPoint (.pptx)
     - Word (.docx)
     - Plain text (.txt)
2. Breaks each file into small chunks of text
3. Saves those chunks into a local database (a folder called "knowledge_base")
   so the chatbot can search through them later

NOTE: Old-style files (.doc, .ppt — without the "x") are NOT supported.
If you have files in those formats, open them in Word/PowerPoint and use
"Save As" to convert them to .docx / .pptx first — takes a few seconds.

You only need to re-run this when you ADD or CHANGE files in "Stats notes".

SETUP:
   pip install pypdf chromadb python-pptx python-docx

RUN:
   python build_knowledge_base.py
"""

import os
import chromadb
from pypdf import PdfReader
from pptx import Presentation
from docx import Document

# ---- SETTINGS ----
NOTES_FOLDER = "Stats notes"       # folder containing your files
DB_FOLDER = "knowledge_base"       # where the "memory" will be saved
CHUNK_SIZE = 800                   # how many characters per chunk (a "piece" of text)
CHUNK_OVERLAP = 100                # slight overlap so we don't cut sentences awkwardly

SUPPORTED_EXTENSIONS = (".pdf", ".pptx", ".docx", ".txt")

# Old formats we recognize but can't read, so we can warn the user clearly
UNSUPPORTED_BUT_KNOWN = (".doc", ".ppt", ".xls", ".xlsx", ".jpg", ".jpeg", ".png")


def read_pdf_text(filepath):
    """Reads a PDF file and returns all its text as one big string."""
    reader = PdfReader(filepath)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text


def read_pptx_text(filepath):
    """Reads a PowerPoint file and returns all text from every slide."""
    prs = Presentation(filepath)
    text = ""
    for slide_number, slide in enumerate(prs.slides, start=1):
        text += f"\n[Slide {slide_number}]\n"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text += shape.text + "\n"
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text:
                            text += cell.text + " "
                    text += "\n"
    return text


def read_docx_text(filepath):
    """Reads a Word document and returns all its paragraph text."""
    doc = Document(filepath)
    text = ""
    for para in doc.paragraphs:
        if para.text:
            text += para.text + "\n"
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text:
                    text += cell.text + " "
            text += "\n"
    return text


def read_txt_text(filepath):
    """Reads a plain text file."""
    # Try common encodings so we don't crash on odd files
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            with open(filepath, "r", encoding=encoding) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    return ""


def read_file_text(filepath):
    """Picks the right reader based on file extension."""
    lower = filepath.lower()
    if lower.endswith(".pdf"):
        return read_pdf_text(filepath)
    elif lower.endswith(".pptx"):
        return read_pptx_text(filepath)
    elif lower.endswith(".docx"):
        return read_docx_text(filepath)
    elif lower.endswith(".txt"):
        return read_txt_text(filepath)
    else:
        return None


def split_into_chunks(text, chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP):
    """
    Splits a long text into smaller overlapping chunks.
    Why overlap? So a sentence that gets cut at a chunk boundary
    still appears in full in the NEXT chunk too.
    """
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return chunks


def main():
    print("=" * 55)
    print("  📚 Building knowledge base from your notes...")
    print("=" * 55)

    if not os.path.isdir(NOTES_FOLDER):
        print(f"\n[Error] Couldn't find a folder called '{NOTES_FOLDER}'.")
        print("Make sure this script is in the same folder as 'Stats notes'.")
        return

    chroma_client = chromadb.PersistentClient(path=DB_FOLDER)
    collection = chroma_client.get_or_create_collection(name="stats_notes")

    all_entries = [
        f for f in os.listdir(NOTES_FOLDER)
        if not f.startswith("~$")  # skip temporary Office lock files
    ]

    supported_files = [f for f in all_entries if f.lower().endswith(SUPPORTED_EXTENSIONS)]
    unsupported_files = [f for f in all_entries if f.lower().endswith(UNSUPPORTED_BUT_KNOWN)]

    if unsupported_files:
        print("\n⚠️  These files were found but can't be read yet:")
        for f in unsupported_files:
            print(f"    - {f}")
        print("    (.doc/.ppt: re-save as .docx/.pptx in Word/PowerPoint first)")
        print("    (images/excel: not supported in this version)\n")

    if not supported_files:
        print(f"\n[Error] No supported files found inside '{NOTES_FOLDER}'.")
        print(f"Supported types: {', '.join(SUPPORTED_EXTENSIONS)}")
        return

    print(f"\nFound {len(supported_files)} supported file(s). Processing...\n")

    chunk_id_counter = 0

    for filename in supported_files:
        filepath = os.path.join(NOTES_FOLDER, filename)
        print(f"  Reading: {filename}")

        try:
            text = read_file_text(filepath)
        except Exception as e:
            print(f"    [Skipped] Couldn't read this file: {e}")
            continue

        if not text or not text.strip():
            print(f"    [Skipped] No readable text found in this file.")
            continue

        chunks = split_into_chunks(text)
        print(f"    Split into {len(chunks)} chunk(s).")

        for chunk in chunks:
            collection.add(
                documents=[chunk],
                metadatas=[{"source": filename}],
                ids=[f"chunk_{chunk_id_counter}"],
            )
            chunk_id_counter += 1

    print(f"\n✅ Done! Saved {chunk_id_counter} chunks to the knowledge base.")
    print(f"   You can now run: python study_assistant.py")


if __name__ == "__main__":
    main()
